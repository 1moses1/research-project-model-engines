"""
Research project course presentation generator.

Creates a Google-Slides-friendly PPTX with embedded PNG charts generated from
the current local result files. Run with an environment that has python-pptx
and Pillow installed, for example:

    "/Users/moiseiradukunda/Documents/CMU/SPRING 2026/TECH STARTUP/pptx_env/bin/python" \
        "Research project Fall to Spring/generate_research_course_presentation.py"
"""

from __future__ import annotations

import json
import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parent
ASSET_DIR = SCRIPT_DIR / "presentation_assets"
OUT_PPTX = SCRIPT_DIR / "AI_Compliance_Auditing_Research_Course_Presentation_Moise.pptx"

RESULTS_JSON = REPO_ROOT / "results" / "audit" / "llm_eval_n200_results.json"
XGB_JSON = REPO_ROOT / "reports" / "xgboost_cv_summary.json"
XGB_FAMILY_CSV = REPO_ROOT / "reports" / "xgboost_per_family_f1.csv"
ADV_JSON = REPO_ROOT / "reports" / "adversarial_validation_v2.json"
LIVE_DASHBOARD = SCRIPT_DIR / "poster_figures" / "live_audit_dashboard.png"


# Slide and color system -------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

C_CMU = RGBColor(0xC4, 0x12, 0x30)
C_NAVY = RGBColor(0x0E, 0x1B, 0x2A)
C_BLUE = RGBColor(0x1F, 0x6F, 0xB2)
C_TEAL = RGBColor(0x0E, 0x8A, 0x8A)
C_GREEN = RGBColor(0x1B, 0x7F, 0x4D)
C_ORANGE = RGBColor(0xE8, 0x77, 0x22)
C_YELLOW = RGBColor(0xF3, 0xB6, 0x1B)
C_RED = RGBColor(0xB8, 0x23, 0x2F)
C_DARK = RGBColor(0x16, 0x1A, 0x23)
C_MUTED = RGBColor(0x5B, 0x66, 0x73)
C_LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
C_LINE = RGBColor(0xD8, 0xDE, 0xE6)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Aptos"


def rgb_tuple(c: RGBColor) -> tuple[int, int, int]:
    return (c[0], c[1], c[2])


def load_json(path: Path, default):
    if path.exists():
        return json.loads(path.read_text())
    return default


def pct(v: float, digits: int = 1) -> str:
    return f"{v:.{digits}f}%"


def safe_float(v, default=0.0) -> float:
    try:
        return float(v)
    except Exception:
        return default


def get_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Helvetica.ttf",
        "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
    ]
    for path in candidates:
        if path and Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, width: int) -> list[str]:
    words = text.split()
    lines, current = [], ""
    for word in words:
        test = word if not current else current + " " + word
        if draw.textbbox((0, 0), test, font=font)[2] <= width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_centered(draw, box, text, font, fill):
    x1, y1, x2, y2 = box
    bb = draw.textbbox((0, 0), text, font=font)
    draw.text((x1 + (x2 - x1 - (bb[2] - bb[0])) / 2, y1 + (y2 - y1 - (bb[3] - bb[1])) / 2), text, font=font, fill=fill)


def chart_bar(
    path: Path,
    title: str,
    labels: list[str],
    values: list[float],
    colors: list[tuple[int, int, int]],
    subtitle: str = "",
    ymax: float = 100,
    ylabel: str = "Percent",
    ci: list[tuple[float, float]] | None = None,
):
    w, h = 1500, 900
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    title_f = get_font(44, True)
    sub_f = get_font(24)
    label_f = get_font(24, True)
    small_f = get_font(22)
    axis_f = get_font(24)

    d.text((70, 45), title, font=title_f, fill=rgb_tuple(C_CMU))
    if subtitle:
        d.text((72, 105), subtitle, font=sub_f, fill=rgb_tuple(C_MUTED))

    left, top, right, bottom = 110, 175, 1440, 730
    d.line((left, bottom, right, bottom), fill=rgb_tuple(C_LINE), width=3)
    d.line((left, top, left, bottom), fill=rgb_tuple(C_LINE), width=3)
    for step in range(0, int(ymax) + 1, 20):
        y = bottom - (step / ymax) * (bottom - top)
        d.line((left, y, right, y), fill=(232, 236, 241), width=1)
        d.text((35, y - 13), str(step), font=axis_f, fill=rgb_tuple(C_MUTED))
    d.text((20, top - 40), ylabel, font=small_f, fill=rgb_tuple(C_MUTED))

    n = len(values)
    gap = 38
    bar_w = (right - left - gap * (n + 1)) / n
    for i, (label, value) in enumerate(zip(labels, values)):
        x1 = left + gap + i * (bar_w + gap)
        x2 = x1 + bar_w
        y1 = bottom - (value / ymax) * (bottom - top)
        d.rounded_rectangle((x1, y1, x2, bottom), radius=18, fill=colors[i])
        d.text((x1 + bar_w / 2 - 36, y1 - 44), pct(value), font=label_f, fill=rgb_tuple(C_DARK))
        if ci:
            lo, hi = ci[i]
            ylo = bottom - (lo / ymax) * (bottom - top)
            yhi = bottom - (hi / ymax) * (bottom - top)
            xc = x1 + bar_w / 2
            d.line((xc, yhi, xc, ylo), fill=rgb_tuple(C_DARK), width=4)
            d.line((xc - 18, yhi, xc + 18, yhi), fill=rgb_tuple(C_DARK), width=4)
            d.line((xc - 18, ylo, xc + 18, ylo), fill=rgb_tuple(C_DARK), width=4)
        lines = label.split("\\n")
        ytxt = bottom + 24
        for line in lines:
            bb = d.textbbox((0, 0), line, font=small_f)
            d.text((x1 + bar_w / 2 - (bb[2] - bb[0]) / 2, ytxt), line, font=small_f, fill=rgb_tuple(C_DARK))
            ytxt += 30
    img.save(path)


def chart_pipeline(path: Path):
    w, h = 1500, 840
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    title_f = get_font(44, True)
    head_f = get_font(26, True)
    body_f = get_font(21)
    small_f = get_font(18)
    d.text((60, 42), "Proposed hybrid compliance-auditing pipeline", font=title_f, fill=rgb_tuple(C_CMU))
    d.text((62, 100), "Fast structured path + semantic LLM path + deterministic control decisions", font=get_font(24), fill=rgb_tuple(C_MUTED))

    boxes = [
        ("Evidence\\nSources", "syslog, Windows Events, HTTP/API, policy docs", (60, 235), rgb_tuple(C_BLUE)),
        ("Normalize +\\nEnrich", "EventID descriptions, fields, taxonomy context", (310, 235), rgb_tuple(C_TEAL)),
        ("Router", "vocabulary coverage + rule confidence", (560, 235), rgb_tuple(C_ORANGE)),
        ("XGBoost v2", "structured / in-vocabulary logs", (810, 160), rgb_tuple(C_GREEN)),
        ("LLM Analyzer", "ambiguous / OOD semantic reasoning", (810, 330), rgb_tuple(C_BLUE)),
        ("Decision\\nEngine", "NCSA thresholds -> compliant / partial / non-compliant", (1060, 235), rgb_tuple(C_CMU)),
        ("Audit\\nReport", "evidence, gaps, remediation", (1310, 235), rgb_tuple(C_NAVY)),
    ]
    for title, body, (x, y), color in boxes:
        d.rounded_rectangle((x, y, x + 170, y + 115), radius=18, fill=color)
        for j, line in enumerate(title.split("\\n")):
            draw_centered(d, (x + 10, y + 12 + j * 24, x + 160, y + 42 + j * 24), line, head_f, "white")
        for j, line in enumerate(wrap_text(d, body, small_f, 150)[:3]):
            draw_centered(d, (x + 10, y + 70 + j * 18, x + 160, y + 95 + j * 18), line, small_f, "white")
    arrows = [(230, 292, 310, 292), (480, 292, 560, 292), (730, 292, 810, 218), (730, 292, 810, 388), (980, 218, 1060, 292), (980, 388, 1060, 292), (1230, 292, 1310, 292)]
    for x1, y1, x2, y2 in arrows:
        d.line((x1, y1, x2, y2), fill=rgb_tuple(C_DARK), width=5)
        angle = math.atan2(y2 - y1, x2 - x1)
        head = 16
        p1 = (x2 - head * math.cos(angle - 0.45), y2 - head * math.sin(angle - 0.45))
        p2 = (x2 - head * math.cos(angle + 0.45), y2 - head * math.sin(angle + 0.45))
        d.polygon([(x2, y2), p1, p2], fill=rgb_tuple(C_DARK))

    d.rounded_rectangle((165, 560, 1335, 735), radius=22, fill=(244, 246, 248), outline=rgb_tuple(C_LINE), width=2)
    d.text((205, 590), "Core research claim", font=head_f, fill=rgb_tuple(C_DARK))
    d.text((205, 635), "XGBoost is efficient but vocabulary-bound; LLM semantic analysis handles heterogeneous logs after normalization.", font=body_f, fill=rgb_tuple(C_MUTED))
    d.text((205, 675), "Decision quality is scored as C(k) = alpha*D_k + beta*V_k + gamma*R_k.", font=body_f, fill=rgb_tuple(C_MUTED))
    img.save(path)


def chart_taxonomy(path: Path):
    w, h = 1500, 860
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    d.text((60, 40), "Control taxonomy enables evidence-to-control learning", font=get_font(44, True), fill=rgb_tuple(C_CMU))
    d.text((62, 100), "Rwanda NCSA controls are mapped to NIST-style families and log evidence patterns.", font=get_font(24), fill=rgb_tuple(C_MUTED))
    families = [
        ("AC", "Access Control", 52, 35, rgb_tuple(C_BLUE)),
        ("AU", "Audit and Accountability", 30, 22, rgb_tuple(C_GREEN)),
        ("CM", "Configuration Mgmt.", 18, 12, rgb_tuple(C_ORANGE)),
        ("IA", "Identity/Auth", 15, 10, rgb_tuple(C_TEAL)),
        ("SC", "System/Comm. Protection", 26, 15, rgb_tuple(C_CMU)),
        ("SI", "System Integrity", 2, 2, rgb_tuple(C_YELLOW)),
    ]
    x0, y0 = 80, 190
    for i, (code, name, total, auditable, color) in enumerate(families):
        x = x0 + (i % 3) * 460
        y = y0 + (i // 3) * 175
        d.rounded_rectangle((x, y, x + 390, y + 120), radius=20, fill=(247, 249, 251), outline=color, width=4)
        d.ellipse((x + 22, y + 24, x + 92, y + 94), fill=color)
        draw_centered(d, (x + 22, y + 24, x + 92, y + 94), code, get_font(24, True), "white")
        d.text((x + 115, y + 22), name, font=get_font(25, True), fill=rgb_tuple(C_DARK))
        d.text((x + 115, y + 62), f"{total} controls | {auditable} system-auditable", font=get_font(21), fill=rgb_tuple(C_MUTED))
    d.rounded_rectangle((150, 580, 1350, 735), radius=24, fill=rgb_tuple(C_NAVY))
    d.text((195, 612), "Training mapping example", font=get_font(27, True), fill="white")
    d.text((195, 660), "Failed SSH login -> AC-7 lockout, IA-2 authentication, AU-2 audit event, SI-4 monitoring", font=get_font(25), fill=(210, 228, 246))
    d.text((195, 700), "One log can activate multiple compliance controls; binary labels lose this structure.", font=get_font(22), fill=(188, 200, 214))
    img.save(path)


def chart_windows_enrichment(path: Path):
    labels = ["Bare\\nEventID", "EventID +\\nDescription"]
    vals = [3.0, 97.5]
    cis = [(1.4, 6.4), (94.3, 98.9)]
    chart_bar(
        path,
        "Windows EventID normalization changes the result",
        labels,
        vals,
        [rgb_tuple(C_RED), rgb_tuple(C_GREEN)],
        "GPT-4o-mini on Mordor Windows Security Events, n=200",
        ymax=100,
        ylabel="Accuracy",
        ci=cis,
    )


def chart_scorecard(path: Path, xgb, llm, adv):
    w, h = 1500, 850
    img = Image.new("RGB", (w, h), "white")
    d = ImageDraw.Draw(img)
    d.text((60, 42), "Current validated results for today's presentation", font=get_font(44, True), fill=rgb_tuple(C_CMU))
    d.text((62, 102), "These numbers are newer than the manuscript/poster abstract and should drive the talk.", font=get_font(24), fill=rgb_tuple(C_MUTED))
    cards = [
        ("Clean XGBoost v2", pct(xgb["test_f1"] * 100, 2), "held-out test F1", rgb_tuple(C_GREEN)),
        ("GPT-4o-mini", pct(llm["macro_accuracy_pct"], 2), "macro accuracy, n=800", rgb_tuple(C_BLUE)),
        ("Windows enrichment", "3.0 -> 97.5%", "bare EventID to described EventID", rgb_tuple(C_TEAL)),
        ("MITRE validation", pct(adv["macro_detection_rate_pct"], 1), "macro detection across 5 scenarios", rgb_tuple(C_ORANGE)),
    ]
    for i, (title, num, label, color) in enumerate(cards):
        x = 90 + i * 350
        y = 230
        d.rounded_rectangle((x, y, x + 300, y + 230), radius=28, fill=(248, 250, 252), outline=color, width=5)
        d.text((x + 25, y + 28), title, font=get_font(26, True), fill=rgb_tuple(C_DARK))
        d.text((x + 25, y + 86), num, font=get_font(42, True), fill=color)
        for j, line in enumerate(wrap_text(d, label, get_font(22), 245)):
            d.text((x + 25, y + 155 + j * 28), line, font=get_font(22), fill=rgb_tuple(C_MUTED))
    d.rounded_rectangle((120, 565, 1380, 720), radius=22, fill=rgb_tuple(C_NAVY))
    d.text((165, 595), "Presentation framing", font=get_font(28, True), fill="white")
    d.text((165, 642), "The strongest contribution is not a perfect model; it is a defensible hybrid architecture that exposes and handles domain shift.", font=get_font(25), fill=(210, 228, 246))
    d.text((165, 682), "Say explicitly: the manuscript is under peer review and the results refresh is being incorporated.", font=get_font(23), fill=(188, 200, 214))
    img.save(path)


def generate_charts(data):
    ASSET_DIR.mkdir(exist_ok=True)

    xgb = data["xgb"]
    llm = data["llm"]
    adv = data["adv"]
    llm_rows = data["llm_rows"]
    adv_rows = data["adv_rows"]
    family_rows = data["family_rows"]

    chart_scorecard(ASSET_DIR / "scorecard.png", xgb, llm, adv)
    chart_pipeline(ASSET_DIR / "pipeline.png")
    chart_taxonomy(ASSET_DIR / "taxonomy.png")

    chart_bar(
        ASSET_DIR / "xgb_clean.png",
        "XGBoost after leakage cleanup",
        ["Leaky\\nold claim", "Clean v2\\nCV F1", "Clean v2\\ntest F1", "Clean v2\\ntest acc."],
        [99.99, xgb["mean_f1"] * 100, xgb["test_f1"] * 100, xgb["test_accuracy"] * 100],
        [rgb_tuple(C_RED), rgb_tuple(C_GREEN), rgb_tuple(C_GREEN), rgb_tuple(C_BLUE)],
        "Clean model is credible: lower performance, but defensible.",
        ymax=110,
        ylabel="Percent",
    )

    labels = [r["log_type"].replace(" Authentication Logs", "\\nAuth").replace(" Access Logs", "\\nAccess").replace("System/Service", "System\\nService").replace(" Security Events", "\\nEvents") for r in llm_rows]
    values = [safe_float(r["accuracy_pct"]) for r in llm_rows]
    ci = [(safe_float(r["wilson_ci_95"][0]), safe_float(r["wilson_ci_95"][1])) for r in llm_rows]
    chart_bar(
        ASSET_DIR / "llm_n800.png",
        "LLM semantic path handles heterogeneous logs",
        labels,
        values,
        [rgb_tuple(C_BLUE), rgb_tuple(C_TEAL), rgb_tuple(C_ORANGE), rgb_tuple(C_GREEN)],
        f"GPT-4o-mini macro accuracy: {pct(llm['macro_accuracy_pct'], 2)}; n={llm['total_n']}",
        ymax=110,
        ylabel="Accuracy",
        ci=ci,
    )

    chart_windows_enrichment(ASSET_DIR / "windows_enrichment.png")

    adv_labels = [r["mitre_id"] + "\\n" + r["scenario"].replace(": ", "\\n")[:28] for r in adv_rows]
    adv_vals = [safe_float(r["detection_rate_pct"]) for r in adv_rows]
    chart_bar(
        ASSET_DIR / "adversarial.png",
        "Effectiveness validation expanded to 5 MITRE scenarios",
        adv_labels,
        adv_vals,
        [rgb_tuple(C_BLUE), rgb_tuple(C_GREEN), rgb_tuple(C_TEAL), rgb_tuple(C_ORANGE), rgb_tuple(C_CMU)],
        f"Macro detection rate: {pct(adv['macro_detection_rate_pct'], 1)}",
        ymax=110,
        ylabel="Detection rate",
    )

    fam_labels = [r["family"].replace(" and ", " &\\n").replace("System ", "System\\n")[:22] for r in family_rows[:8]]
    fam_vals = [safe_float(r["f1_binary"]) * 100 for r in family_rows[:8]]
    chart_bar(
        ASSET_DIR / "family_f1.png",
        "Clean XGBoost performance varies by control family",
        fam_labels,
        fam_vals,
        [rgb_tuple(C_BLUE), rgb_tuple(C_GREEN), rgb_tuple(C_ORANGE), rgb_tuple(C_TEAL), rgb_tuple(C_CMU), rgb_tuple(C_YELLOW), rgb_tuple(C_NAVY), rgb_tuple(C_RED)],
        "Per-family F1 from leakage-free held-out test set",
        ymax=110,
        ylabel="F1",
    )


def parse_family_csv(path: Path):
    if not path.exists():
        return []
    rows = []
    lines = path.read_text().strip().splitlines()
    headers = lines[0].split(",")
    for line in lines[1:]:
        vals = line.split(",")
        rows.append(dict(zip(headers, vals)))
    rows.sort(key=lambda r: safe_float(r.get("f1_binary")), reverse=True)
    return rows


def load_data():
    xgb_raw = load_json(XGB_JSON, {})
    xgb = {
        "mean_f1": safe_float(xgb_raw.get("mean_f1", 0.852)),
        "test_f1": safe_float(xgb_raw.get("test_f1", 0.8545)),
        "test_accuracy": safe_float(xgb_raw.get("test_accuracy", 0.8714)),
        "ci": xgb_raw.get("wilson_ci_95", [0.8491, 0.8549]),
    }

    llm_raw = load_json(RESULTS_JSON, {})
    rows = llm_raw.get("results", [])
    if not rows:
        rows = [
            {"log_type": "SSH Authentication Logs", "accuracy_pct": 89.5, "wilson_ci_95": [84.5, 93.0], "correct": 179, "n_samples": 200},
            {"log_type": "Windows Security Events", "accuracy_pct": 97.5, "wilson_ci_95": [94.3, 98.9], "correct": 195, "n_samples": 200},
            {"log_type": "HTTP/API Access Logs", "accuracy_pct": 95.5, "wilson_ci_95": [91.7, 97.6], "correct": 191, "n_samples": 200},
            {"log_type": "macOS System/Service Logs", "accuracy_pct": 96.5, "wilson_ci_95": [93.0, 98.3], "correct": 193, "n_samples": 200},
        ]
    total_correct = sum(int(r.get("correct", 0)) for r in rows)
    total_n = sum(int(r.get("n_samples", 0)) for r in rows)
    llm = llm_raw.get("overall", {})
    llm = {
        "total_correct": int(llm.get("total_correct", total_correct)),
        "total_n": int(llm.get("total_n", total_n)),
        "macro_accuracy_pct": safe_float(llm.get("macro_accuracy_pct", total_correct / max(total_n, 1) * 100)),
        "wilson_ci_95": llm.get("wilson_ci_95", [92.8, 96.2]),
    }

    adv_raw = load_json(ADV_JSON, {})
    adv_rows = adv_raw.get("scenarios", [])
    adv = {
        "macro_detection_rate_pct": safe_float(adv_raw.get("macro_detection_rate_pct", 92.8)),
        "baseline_fpr": safe_float(adv_raw.get("compliant_baseline", {}).get("false_positive_rate_pct", 53.0)),
    }

    return {
        "xgb": xgb,
        "llm": llm,
        "llm_rows": rows,
        "adv": adv,
        "adv_rows": adv_rows,
        "family_rows": parse_family_csv(XGB_FAMILY_CSV),
    }


# PPTX helpers ----------------------------------------------------------------
def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=C_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill, border=None, radius=False):
    shape_id = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_id, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1.1)
    else:
        sh.line.fill.background()
    return sh


def shape_text(shape, title, body="", title_size=15, body_size=9.5, color=C_WHITE):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = color
    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(5)
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = FONT
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = color


def flow_line(slide, x1, y1, x2, y2, color=C_DARK):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.25)
    return line


def tb(slide, text, left, top, width, height, size=18, color=C_DARK, bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return box


def header(slide, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.92), C_NAVY)
    rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, C_CMU)
    tb(slide, title, Inches(0.45), Inches(0.12), Inches(10.5), Inches(0.45), 24, C_WHITE, True)
    if subtitle:
        tb(slide, subtitle, Inches(0.46), Inches(0.58), Inches(10.8), Inches(0.25), 10.5, RGBColor(0xB9, 0xC7, 0xD8))


def footer(slide, n, total):
    tb(slide, "AI-Augmented Compliance Auditing | Research Project Course | Moise Iradukunda Ingabire", Inches(0.45), Inches(7.12), Inches(9.2), Inches(0.22), 8.5, C_MUTED)
    tb(slide, f"{n}/{total}", Inches(12.25), Inches(7.08), Inches(0.7), Inches(0.22), 9, C_MUTED, align=PP_ALIGN.RIGHT)


def picture(slide, path, left, top, width=None, height=None):
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def stat_card(slide, title, number, caption, left, top, color):
    rect(slide, left, top, Inches(2.85), Inches(1.45), RGBColor(0xF8, 0xFA, 0xFC), color, True)
    tb(slide, title, left + Inches(0.16), top + Inches(0.12), Inches(2.55), Inches(0.25), 11, C_DARK, True)
    tb(slide, number, left + Inches(0.16), top + Inches(0.42), Inches(2.55), Inches(0.45), 24, color, True)
    tb(slide, caption, left + Inches(0.16), top + Inches(0.92), Inches(2.55), Inches(0.35), 9.5, C_MUTED)


def add_table(slide, rows, headers, left, top, width, height, font_size=10):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    col_w = width / len(headers)
    for i in range(len(headers)):
        tbl.columns[i].width = int(col_w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = C_WHITE
                r.font.bold = True
                r.font.size = Pt(font_size)
                r.font.name = FONT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFB) if i % 2 else C_WHITE
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = C_DARK
                    r.font.size = Pt(font_size)
                    r.font.name = FONT
    return tbl


def editable_bar_chart(slide, title, subtitle, labels, values, colors, left, top, width, height, ymax=100, ci=None):
    tb(slide, title, left, top, width, Inches(0.32), 17, C_CMU, True)
    if subtitle:
        tb(slide, subtitle, left, top + Inches(0.36), width, Inches(0.25), 10.5, C_MUTED)

    plot_l = left + Inches(0.45)
    plot_t = top + Inches(0.88)
    plot_w = width - Inches(0.65)
    plot_h = height - Inches(1.55)
    base_y = plot_t + plot_h
    rect(slide, plot_l, base_y, plot_w, Pt(1), C_LINE)
    rect(slide, plot_l, plot_t, Pt(1), plot_h, C_LINE)

    for tick in range(0, int(ymax) + 1, 20):
        y = base_y - plot_h * tick / ymax
        rect(slide, plot_l, y, plot_w, Pt(0.35), RGBColor(0xEC, 0xF0, 0xF4))
        tb(slide, str(tick), left, y - Inches(0.08), Inches(0.35), Inches(0.18), 7.5, C_MUTED, align=PP_ALIGN.RIGHT)

    n = len(values)
    gap = Inches(0.15)
    bar_w = int((plot_w - gap * (n + 1)) / n)
    for i, (label, value) in enumerate(zip(labels, values)):
        x = plot_l + gap + i * (bar_w + gap)
        bar_h = plot_h * value / ymax
        y = base_y - bar_h
        rect(slide, x, y, bar_w, bar_h, colors[i], None, True)
        tb(slide, pct(value), x - Inches(0.06), y - Inches(0.28), bar_w + Inches(0.12), Inches(0.22), 9.5, C_DARK, True, align=PP_ALIGN.CENTER)
        if ci:
            lo, hi = ci[i]
            ylo = base_y - plot_h * lo / ymax
            yhi = base_y - plot_h * hi / ymax
            xc = x + bar_w / 2
            rect(slide, xc - Pt(0.7), yhi, Pt(1.4), ylo - yhi, C_DARK)
            rect(slide, xc - Inches(0.06), yhi, Inches(0.12), Pt(1.2), C_DARK)
            rect(slide, xc - Inches(0.06), ylo, Inches(0.12), Pt(1.2), C_DARK)
        tb(slide, label.replace("\\n", "\n"), x - Inches(0.08), base_y + Inches(0.08), bar_w + Inches(0.16), Inches(0.6), 8.2, C_DARK, align=PP_ALIGN.CENTER)


def editable_scorecard(slide, data):
    cards = [
        ("Clean XGBoost v2", pct(data["xgb"]["test_f1"] * 100, 2), "held-out test F1", C_GREEN),
        ("GPT-4o-mini", pct(data["llm"]["macro_accuracy_pct"], 2), "macro accuracy, n=800", C_BLUE),
        ("Windows enrichment", "3.0 -> 97.5%", "bare EventID to description", C_TEAL),
        ("MITRE validation", pct(data["adv"]["macro_detection_rate_pct"], 1), "5-scenario macro detection", C_ORANGE),
    ]
    for i, (title, number, caption, color) in enumerate(cards):
        x = Inches(0.65 + i * 3.15)
        y = Inches(1.45)
        rect(slide, x, y, Inches(2.85), Inches(1.7), RGBColor(0xF8, 0xFA, 0xFC), color, True)
        tb(slide, title, x + Inches(0.15), y + Inches(0.12), Inches(2.55), Inches(0.25), 11, C_DARK, True)
        tb(slide, number, x + Inches(0.15), y + Inches(0.48), Inches(2.55), Inches(0.45), 24, color, True)
        tb(slide, caption, x + Inches(0.15), y + Inches(1.12), Inches(2.55), Inches(0.35), 9.5, C_MUTED)
    rect(slide, Inches(1.0), Inches(4.35), Inches(11.3), Inches(1.55), C_NAVY, None, True)
    tb(slide, "Presentation framing", Inches(1.25), Inches(4.62), Inches(3.0), Inches(0.3), 17, C_WHITE, True)
    tb(slide, "The contribution is not a perfect model. It is a defensible hybrid architecture that exposes domain shift, normalizes evidence, and routes logs to the right analysis path.", Inches(1.25), Inches(5.08), Inches(10.6), Inches(0.42), 14, RGBColor(0xD2, 0xE5, 0xF6))


def editable_taxonomy(slide):
    families = [
        ("AC", "Access Control", 52, 35, C_BLUE),
        ("AU", "Audit & Accountability", 30, 22, C_GREEN),
        ("CM", "Configuration Mgmt.", 18, 12, C_ORANGE),
        ("IA", "Identity/Auth", 15, 10, C_TEAL),
        ("SC", "System/Comm. Protection", 26, 15, C_CMU),
        ("SI", "System Integrity", 2, 2, C_YELLOW),
    ]
    for i, (code, name, total, auditable, color) in enumerate(families):
        x = Inches(0.75 + (i % 3) * 4.18)
        y = Inches(1.35 + (i // 3) * 1.35)
        rect(slide, x, y, Inches(3.55), Inches(1.0), RGBColor(0xF8, 0xFA, 0xFC), color, True)
        rect(slide, x + Inches(0.18), y + Inches(0.22), Inches(0.55), Inches(0.55), color, None, True)
        tb(slide, code, x + Inches(0.22), y + Inches(0.34), Inches(0.47), Inches(0.18), 12, C_WHITE, True, align=PP_ALIGN.CENTER)
        tb(slide, name, x + Inches(0.9), y + Inches(0.18), Inches(2.45), Inches(0.25), 13.2, C_DARK, True)
        tb(slide, f"{total} controls | {auditable} system-auditable", x + Inches(0.9), y + Inches(0.58), Inches(2.45), Inches(0.25), 10, C_MUTED)
    rect(slide, Inches(1.1), Inches(4.5), Inches(11.1), Inches(1.35), C_NAVY, None, True)
    tb(slide, "Training mapping example", Inches(1.35), Inches(4.74), Inches(4.1), Inches(0.25), 16, C_WHITE, True)
    tb(slide, "Failed SSH login -> AC-7 lockout, IA-2 authentication, AU-2 audit event, SI-4 monitoring", Inches(1.35), Inches(5.18), Inches(10.3), Inches(0.28), 13.5, RGBColor(0xD2, 0xE5, 0xF6))
    tb(slide, "One log can activate multiple controls; binary labels lose compliance structure.", Inches(1.35), Inches(5.52), Inches(10.3), Inches(0.25), 11.5, RGBColor(0xB9, 0xC7, 0xD8))


def editable_pipeline(slide):
    boxes = [
        ("Evidence\nSources", "syslog / Windows\nHTTP / docs", 0.55, 2.45, 1.6, 1.15, C_BLUE),
        ("Normalize +\nEnrich", "EventID -> meaning\nfields -> context", 2.55, 2.45, 1.6, 1.15, C_TEAL),
        ("Router", "coverage +\nconfidence", 4.55, 2.45, 1.55, 1.15, C_ORANGE),
        ("XGBoost v2", "structured\nin-vocab logs", 6.75, 1.65, 1.55, 1.1, C_GREEN),
        ("LLM Analyzer", "ambiguous / OOD\nsemantic logs", 6.75, 3.35, 1.55, 1.1, C_BLUE),
        ("Decision\nEngine", "NCSA thresholds\n3-way status", 9.0, 2.45, 1.65, 1.15, C_CMU),
        ("Audit\nReport", "evidence, gaps\nremediation", 11.25, 2.45, 1.45, 1.15, C_NAVY),
    ]
    centers = {}
    for title, body, x, y, w, h, color in boxes:
        sh = rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), color, None, True)
        shape_text(sh, title, body, 12.8 if w < 1.6 else 13.6, 8.2)
        centers[title] = (Inches(x + w / 2), Inches(y + h / 2))
    def c(name):
        return centers[name]
    flow_line(slide, Inches(2.15), c("Evidence\nSources")[1], Inches(2.55), c("Normalize +\nEnrich")[1])
    flow_line(slide, Inches(4.15), c("Normalize +\nEnrich")[1], Inches(4.55), c("Router")[1])
    flow_line(slide, Inches(6.1), c("Router")[1], Inches(6.75), c("XGBoost v2")[1])
    flow_line(slide, Inches(6.1), c("Router")[1], Inches(6.75), c("LLM Analyzer")[1])
    flow_line(slide, Inches(8.3), c("XGBoost v2")[1], Inches(9.0), c("Decision\nEngine")[1])
    flow_line(slide, Inches(8.3), c("LLM Analyzer")[1], Inches(9.0), c("Decision\nEngine")[1])
    flow_line(slide, Inches(10.65), c("Decision\nEngine")[1], Inches(11.25), c("Audit\nReport")[1])

    rect(slide, Inches(1.15), Inches(5.2), Inches(11.0), Inches(1.05), RGBColor(0xF4, 0xF6, 0xF8), C_LINE, True)
    tb(slide, "Core research claim", Inches(1.45), Inches(5.42), Inches(3.1), Inches(0.25), 15, C_DARK, True)
    tb(slide, "XGBoost is efficient but vocabulary-bound; LLM semantic analysis handles heterogeneous logs after normalization.", Inches(1.45), Inches(5.82), Inches(9.8), Inches(0.22), 11.5, C_MUTED)


def slides(prs, data):
    total = 15
    n = 1

    # 1 Cover
    s = blank(prs); bg(s, C_NAVY)
    rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), C_CMU)
    rect(s, Inches(0.72), Inches(0.78), Inches(0.08), Inches(4.8), C_CMU)
    tb(s, "AI-Augmented\nCompliance Auditing\nfor Cloud Systems", Inches(1.05), Inches(0.82), Inches(10.8), Inches(2.35), 40, C_WHITE, True)
    tb(s, "A hybrid ML-LLM approach for Rwanda NCSA controls", Inches(1.08), Inches(3.25), Inches(10.8), Inches(0.45), 19, RGBColor(0xB9, 0xD9, 0xF2))
    tb(s, "Research Project Course Presentation | May 06, 2026", Inches(1.08), Inches(4.05), Inches(9.5), Inches(0.3), 13, RGBColor(0xC9, 0xD4, 0xDF))
    tb(s, "Moise Iradukunda Ingabire\nMSIT, Carnegie Mellon University Africa", Inches(1.08), Inches(5.55), Inches(6.2), Inches(0.62), 15, C_WHITE, True)
    tb(s, "Supervisor: Prof. Jema David Ndibwile\nManuscript status: submitted to Future Internet; under peer review / major revision refresh", Inches(7.1), Inches(5.55), Inches(5.4), Inches(0.72), 12.5, RGBColor(0xDD, 0xE8, 0xF4))
    tb(s, str(n), Inches(12.45), Inches(7.02), Inches(0.5), Inches(0.25), 10, RGBColor(0xC9, 0xD4, 0xDF), align=PP_ALIGN.RIGHT)
    n += 1

    # 2 Storyline
    s = blank(prs); bg(s); header(s, "Today's story", "Keep the talk anchored on the new validated results, not the stale abstract numbers.")
    cards = [
        ("1", "Problem", "Audits verify presence, not whether controls work."),
        ("2", "Gap", "Vocabulary-based ML collapses across log formats."),
        ("3", "Solution", "Hybrid XGBoost + LLM + deterministic decision engine."),
        ("4", "Evidence", "Clean XGBoost, n=800 LLM eval, 5 MITRE scenarios."),
    ]
    for i, (idx, title, desc) in enumerate(cards):
        x = Inches(0.7 + i * 3.1)
        rect(s, x, Inches(1.55), Inches(2.75), Inches(3.6), RGBColor(0xF8, 0xFA, 0xFC), C_LINE, True)
        rect(s, x + Inches(0.2), Inches(1.82), Inches(0.55), Inches(0.55), [C_BLUE, C_TEAL, C_ORANGE, C_CMU][i], None, True)
        tb(s, idx, x + Inches(0.33), Inches(1.93), Inches(0.3), Inches(0.25), 15, C_WHITE, True, align=PP_ALIGN.CENTER)
        tb(s, title, x + Inches(0.22), Inches(2.62), Inches(2.3), Inches(0.35), 18, C_DARK, True)
        tb(s, desc, x + Inches(0.22), Inches(3.12), Inches(2.25), Inches(1.1), 13, C_MUTED)
    tb(s, "One-line thesis: The contribution is an honest, low-cost hybrid architecture that detects domain shift and routes evidence to the right analysis path.", Inches(1.0), Inches(5.78), Inches(11.2), Inches(0.6), 18, C_NAVY, True, align=PP_ALIGN.CENTER)
    footer(s, n, total); n += 1

    # 3 Problem
    s = blank(prs); bg(s); header(s, "Research problem", "Manual compliance auditing is slow, expensive, and often checkbox-driven.")
    stat_card(s, "Audit burden", "1000+ h", "reported annual manual audit effort", Inches(0.7), Inches(1.35), C_CMU)
    stat_card(s, "Budget pressure", "40%", "IT security budget consumed by compliance", Inches(3.65), Inches(1.35), C_BLUE)
    stat_card(s, "Control scope", "143", "Rwanda NCSA controls modeled", Inches(6.6), Inches(1.35), C_GREEN)
    stat_card(s, "Core flaw", "Presence != Effectiveness", "installed controls can still fail", Inches(9.55), Inches(1.35), C_ORANGE)
    tb(s, "Why this matters", Inches(0.8), Inches(3.35), Inches(3.4), Inches(0.35), 20, C_DARK, True)
    tb(s, "A firewall, WAF, or EDR can be present but fail against real attack behavior. The research asks how to move from checkbox compliance to measured control effectiveness.", Inches(0.8), Inches(3.85), Inches(5.2), Inches(1.6), 15, C_MUTED)
    tb(s, "Research question", Inches(7.0), Inches(3.35), Inches(3.4), Inches(0.35), 20, C_DARK, True)
    tb(s, "Can a low-cost hybrid ML-LLM auditor map heterogeneous evidence to Rwanda NCSA controls and expose effectiveness gaps under realistic logs and adversarial tests?", Inches(7.0), Inches(3.85), Inches(5.15), Inches(1.6), 15, C_MUTED)
    footer(s, n, total); n += 1

    # 4 Literature gap
    s = blank(prs); bg(s); header(s, "Literature review: what is missing?", "Existing tools solve pieces; few combine evidence mapping, semantic logs, and adversarial effectiveness.")
    rows = [
        ["OpenSCAP / CIS-CAT", "Rule checks", "No semantic/OOD log reasoning"],
        ["Wazuh / SIEM rules", "Operational detection", "Limited regulatory effectiveness scoring"],
        ["BERT / LSTM log ML", "High accuracy in-domain", "GPU and domain-shift costs"],
        ["Commercial GRC", "Workflow management", "High cost; little adversarial validation"],
        ["This research", "Hybrid ML-LLM + NCSA mapping", "Still needs institutional pilot validation"],
    ]
    add_table(s, rows, ["Approach", "Strength", "Gap addressed"], Inches(0.7), Inches(1.25), Inches(12.0), Inches(3.8), 11)
    tb(s, "Positioning", Inches(1.0), Inches(5.4), Inches(11.0), Inches(0.75), 18, C_NAVY, True, align=PP_ALIGN.CENTER)
    tb(s, "The work is strongest when framed as a reference architecture and validation pipeline for resource-constrained national compliance programs, piloting from Rwanda NCSA.", Inches(1.2), Inches(6.03), Inches(10.6), Inches(0.5), 14, C_MUTED, align=PP_ALIGN.CENTER)
    footer(s, n, total); n += 1

    # 5 Taxonomy
    s = blank(prs); bg(s); header(s, "Control taxonomy and evidence mapping", "The taxonomy is the bridge between raw logs and trainable compliance labels.")
    editable_taxonomy(s)
    footer(s, n, total); n += 1

    # 6 Architecture
    s = blank(prs); bg(s); header(s, "Proposed research solution", "Use the right analysis path for the right evidence type.")
    editable_pipeline(s)
    footer(s, n, total); n += 1

    # 7 Current results scorecard
    s = blank(prs); bg(s); header(s, "Current validated results", "These are the numbers to present today.")
    editable_scorecard(s, data)
    footer(s, n, total); n += 1

    # 8 Data integrity refresh
    s = blank(prs); bg(s); header(s, "Peer review forced a stronger result", "The old near-perfect metric was not defensible; the clean model is.")
    editable_bar_chart(
        s,
        "XGBoost after leakage cleanup",
        "Clean model is credible: lower performance, but defensible.",
        ["Leaky\nold claim", "Clean v2\nCV F1", "Clean v2\ntest F1", "Clean v2\ntest acc."],
        [99.99, data["xgb"]["mean_f1"] * 100, data["xgb"]["test_f1"] * 100, data["xgb"]["test_accuracy"] * 100],
        [C_RED, C_GREEN, C_GREEN, C_BLUE],
        Inches(0.65), Inches(1.15), Inches(7.05), Inches(4.65), ymax=110,
    )
    rows = [
        ["Old issue", "Synthetic leakage inflated the headline F1"],
        ["Fix", "Regenerated leakage-free data; removed leaky features"],
        ["New result", f"Clean XGBoost v2 test F1 = {pct(data['xgb']['test_f1'] * 100, 2)}"],
        ["Interpretation", "Credible baseline in the expected 70-90% range"],
    ]
    add_table(s, rows, ["Item", "Status"], Inches(8.2), Inches(1.55), Inches(4.35), Inches(3.65), 11)
    tb(s, "How to say it: reviewers were right to challenge the perfect score; the revision now reports the honest model behavior.", Inches(8.25), Inches(5.62), Inches(4.25), Inches(0.75), 13.5, C_NAVY, True)
    footer(s, n, total); n += 1

    # 9 Per-family XGB
    s = blank(prs); bg(s); header(s, "XGBoost still has a role", "It is fast and useful for structured, in-distribution evidence.")
    fam_rows = data["family_rows"][:8]
    family_code = {
        "Access Control": "AC",
        "Audit and Accountability": "AU",
        "Awareness and Training": "AT",
        "Configuration Management": "CM",
        "Identity Management and Authentication": "IA",
        "Incident Response": "IR",
        "Maintenance": "MA",
        "Media Protection": "MP",
        "Personnel Security": "PS",
        "Physical and Environmental Protection": "PE",
        "Risk Assessment": "RA",
        "Security Assessment": "SA",
        "Security Policy and Procedures": "SP",
        "System and Communications Protection": "SC",
    }
    editable_bar_chart(
        s,
        "Performance varies by control family",
        "Per-family F1 from leakage-free held-out test set.",
        [family_code.get(r["family"], r["family"][:4]) for r in fam_rows],
        [safe_float(r["f1_binary"]) * 100 for r in fam_rows],
        [C_BLUE, C_GREEN, C_ORANGE, C_TEAL, C_CMU, C_YELLOW, C_NAVY, C_RED],
        Inches(0.65), Inches(1.15), Inches(7.55), Inches(4.95), ymax=110,
    )
    tb(s, "What the chart means", Inches(8.55), Inches(1.35), Inches(3.5), Inches(0.35), 20, C_DARK, True)
    tb(s, "The clean XGBoost model is not the universal compliance brain. It is the low-latency fast path for log families with vocabulary overlap and structured evidence.", Inches(8.55), Inches(1.92), Inches(3.85), Inches(1.35), 15, C_MUTED)
    tb(s, "Design decision", Inches(8.55), Inches(3.65), Inches(3.5), Inches(0.35), 20, C_DARK, True)
    tb(s, "When vocabulary coverage is weak, route to the semantic LLM path before the Decision Engine makes the control verdict.", Inches(8.55), Inches(4.2), Inches(3.85), Inches(1.1), 15, C_MUTED)
    footer(s, n, total); n += 1

    # 10 LLM
    s = blank(prs); bg(s); header(s, "LLM semantic analysis closes the format gap", "Current GPT-4o-mini evaluation: 800 samples, 4 log types, Wilson confidence intervals.")
    llm_rows = data["llm_rows"]
    editable_bar_chart(
        s,
        "LLM semantic path handles heterogeneous logs",
        f"GPT-4o-mini macro accuracy: {pct(data['llm']['macro_accuracy_pct'], 2)}; n={data['llm']['total_n']}",
        [r["log_type"].replace(" Authentication Logs", "\nAuth").replace(" Access Logs", "\nAccess").replace("System/Service", "System\nService").replace(" Security Events", "\nEvents") for r in llm_rows],
        [safe_float(r["accuracy_pct"]) for r in llm_rows],
        [C_BLUE, C_TEAL, C_ORANGE, C_GREEN],
        Inches(0.75), Inches(1.15), Inches(11.8), Inches(5.25), ymax=110,
        ci=[(safe_float(r["wilson_ci_95"][0]), safe_float(r["wilson_ci_95"][1])) for r in llm_rows],
    )
    footer(s, n, total); n += 1

    # 11 Windows enrichment
    s = blank(prs); bg(s); header(s, "Key architectural finding: normalize before reasoning", "Bare EventID integers are not enough; descriptions unlock semantic classification.")
    editable_bar_chart(
        s,
        "Windows EventID normalization changes the result",
        "GPT-4o-mini on Mordor Windows Security Events, n=200",
        ["Bare\nEventID", "EventID +\nDescription"],
        [3.0, 97.5],
        [C_RED, C_GREEN],
        Inches(0.7), Inches(1.15), Inches(7.15), Inches(4.8), ymax=100,
        ci=[(1.4, 6.4), (94.3, 98.9)],
    )
    tb(s, "Presentation point", Inches(8.45), Inches(1.45), Inches(3.7), Inches(0.35), 20, C_DARK, True)
    tb(s, "Do not claim the LLM magically understands every raw enterprise log. The system must enrich Windows EventIDs into natural-language security context before LLM classification.", Inches(8.45), Inches(2.0), Inches(3.9), Inches(1.45), 15, C_MUTED)
    tb(s, "This supports the hybrid design: normalization is not preprocessing trivia; it is part of the compliance reasoning architecture.", Inches(8.45), Inches(4.15), Inches(3.9), Inches(1.15), 15, C_NAVY, True)
    footer(s, n, total); n += 1

    # 12 Adversarial
    s = blank(prs); bg(s); header(s, "Effectiveness validation: 5 MITRE ATT&CK scenarios", "The system tests whether controls work, not only whether they exist.")
    editable_bar_chart(
        s,
        "Expanded adversarial validation",
        f"Macro detection rate: {pct(data['adv']['macro_detection_rate_pct'], 1)}",
        [r.get("mitre_id", "") for r in data["adv_rows"]],
        [safe_float(r.get("detection_rate_pct")) for r in data["adv_rows"]],
        [C_BLUE, C_GREEN, C_TEAL, C_ORANGE, C_CMU],
        Inches(0.65), Inches(1.15), Inches(7.95), Inches(4.95), ymax=110,
    )
    rows = []
    for r in data["adv_rows"]:
        rows.append([r.get("mitre_id", ""), r.get("scenario", "")[:28], pct(safe_float(r.get("detection_rate_pct")), 1)])
    add_table(s, rows, ["MITRE", "Scenario", "Detection"], Inches(8.9), Inches(1.35), Inches(3.75), Inches(4.35), 9.3)
    tb(s, "Framing: this turns compliance into measured control effectiveness under attack behavior.", Inches(8.9), Inches(6.02), Inches(3.65), Inches(0.45), 12.5, C_NAVY, True)
    footer(s, n, total); n += 1

    # 13 Dataset provenance
    s = blank(prs); bg(s); header(s, "Dataset provenance for the latest results", "Use this slide to explain exactly what evidence supports each headline number.")
    rows = [
        ["Clean XGBoost v2", "Leakage-free regenerated synthetic NCSA logs", "70K train / 15K val / 15K test", "85.45% test F1"],
        ["LLM semantic eval", "SSH auth, Mordor Windows, HTTP/API, macOS logs", "200/type = 800 total", "94.75% GPT-4o-mini"],
        ["Windows enrichment", "Mordor APT3 Windows Security EventIDs", "n=200", "3.0% bare -> 97.5% enriched"],
        ["Adversarial validation", "MITRE T1059.001, T1190, T1110, T1078, T1562.001", "~50/scenario", "92.8% macro detection"],
        ["Compliant FPR check", "Mordor + Elastic AWS auth + Elastic Apache + LogPai Linux", "100 real compliant logs", "53.0% FPR currently"],
    ]
    add_table(
        s,
        rows,
        ["Result", "Dataset / source", "Sample size", "Current number"],
        Inches(0.55), Inches(1.28), Inches(12.25), Inches(4.1),
        9.6,
    )
    rect(s, Inches(0.9), Inches(5.85), Inches(11.65), Inches(0.82), RGBColor(0xF4, 0xF6, 0xF8), C_LINE, True)
    tb(s, "Important wording", Inches(1.15), Inches(6.04), Inches(2.4), Inches(0.24), 13.5, C_DARK, True)
    tb(s, "Do not say one dataset produced every result. The honest story is a pipeline evaluation: clean synthetic training for XGBoost, mixed real/generated evaluation for LLM and adversarial tests, and real compliant logs for false-positive analysis.", Inches(3.45), Inches(5.98), Inches(8.6), Inches(0.42), 11.2, C_MUTED)
    footer(s, n, total); n += 1

    # 14 Deployment/status
    s = blank(prs); bg(s); header(s, "Deployment evidence and research status", "Prototype works; paper is being refreshed for peer review.")
    stat_card(s, "Live audit", "0.77 s", "end-to-end report generation", Inches(0.75), Inches(1.28), C_GREEN)
    stat_card(s, "Footprint", "2 CPU / 2.66 GB", "measured seven-engine system", Inches(3.75), Inches(1.28), C_BLUE)
    stat_card(s, "Hosting target", "~$50/mo + API", "resource-constrained deployment", Inches(6.75), Inches(1.28), C_ORANGE)
    stat_card(s, "Manuscript", "Future Internet", "submitted; under peer review refresh", Inches(9.75), Inches(1.28), C_CMU)
    if LIVE_DASHBOARD.exists():
        picture(s, LIVE_DASHBOARD, Inches(0.9), Inches(3.05), width=Inches(5.8))
    tb(s, "Status for the audience", Inches(7.3), Inches(3.22), Inches(4.8), Inches(0.35), 20, C_DARK, True)
    tb(s, "The manuscript/poster abstract still contains older numbers. Today's presentation uses the current local reports: clean XGBoost v2, n=800 GPT-4o-mini evaluation, Windows enrichment, and 5-scenario MITRE validation.", Inches(7.3), Inches(3.85), Inches(4.9), Inches(1.35), 15, C_MUTED)
    tb(s, "This is a research-in-progress revision, not a production certification claim.", Inches(7.3), Inches(5.55), Inches(4.9), Inches(0.55), 15, C_NAVY, True)
    footer(s, n, total); n += 1

    # 15 Future
    s = blank(prs); bg(s); header(s, "Future improvements", "From Rwanda NCSA pilot to model-agnostic African compliance infrastructure.")
    rows = [
        ["Short term", "Refresh manuscript results; add supplementary tables; finish reviewer responses"],
        ["Pilot", "Validate on production logs from 2-3 Rwandan institutions"],
        ["Frontier models", "Provider-agnostic adapters for Claude Opus, Gemini, GPT-class, Mistral, Mythos-style models"],
        ["On-premise", "Fine-tune Llama/Mistral-family models for data-sovereignty settings"],
        ["Pan-African path", "Map Rwanda NCSA to ISO 27001, NIST SP 800-53, Kenya, Nigeria, and Malabo-aligned frameworks"],
    ]
    add_table(s, rows, ["Horizon", "Work planned"], Inches(0.75), Inches(1.25), Inches(11.85), Inches(4.15), 11)
    tb(s, "Closing line", Inches(1.0), Inches(5.95), Inches(11.2), Inches(0.35), 18, C_DARK, True, align=PP_ALIGN.CENTER)
    tb(s, "The long-term vision is not one model or one country: it is a modular compliance-auditing reference architecture for resource-constrained African institutions.", Inches(1.05), Inches(6.38), Inches(11.1), Inches(0.45), 15, C_NAVY, True, align=PP_ALIGN.CENTER)
    footer(s, n, total)


def main():
    data = load_data()
    generate_charts(data)
    prs = new_prs()
    slides(prs, data)
    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote generator {Path(__file__).resolve()}")


if __name__ == "__main__":
    main()
